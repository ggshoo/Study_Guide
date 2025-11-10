import os
from pathlib import Path
import chromadb
from chromadb.config import Settings
import openai
from langchain.text_splitter import RecursiveCharacterTextSplitter
from tqdm import tqdm
from dotenv import load_dotenv
import PyPDF2
from pptx import Presentation
import re

# Load environment variables
load_dotenv()


class AIStudyAssistant:
    def __init__(self, persist_directory="db"):
        """Initialize the AI Study Assistant with a ChromaDB instance."""
        self.client = chromadb.PersistentClient(path=persist_directory)
        self.collection = self.client.get_or_create_collection(
            name="lecture_content",
            metadata={"hnsw:space": "cosine"}
        )

        # Ensure OpenAI API key is configured
        openai.api_key = os.getenv("OPENAI_API_KEY")
        if not openai.api_key:
            raise ValueError("OPENAI_API_KEY environment variable is not set")

        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )

    def get_embedding(self, text: str):
        """Call OpenAI Embeddings API directly and return the vector."""
        resp = openai.Embedding.create(model="text-embedding-3-small", input=text)
        return resp["data"][0]["embedding"]

    def extract_pptx_content(self, pptx_path: str):
        """Extract text content from PowerPoint file with slide numbers."""
        prs = Presentation(pptx_path)
        slides_content = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = f"\n--- Slide {slide_num} ---\n"
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            
            # Extract notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    slide_text += f"\nNotes: {notes_text}\n"
            
            slides_content.append({
                'slide_number': slide_num,
                'content': slide_text,
                'filename': Path(pptx_path).name
            })
        
        return slides_content

    def process_pptx(self, pptx_path: str):
        """Process PowerPoint file and add to vector database with slide tracking."""
        slides = self.extract_pptx_content(pptx_path)
        
        for slide_info in tqdm(slides, desc=f"Processing {Path(pptx_path).name}"):
            # Don't chunk slides - keep each slide as one unit for accurate slide number tracking
            embedding = self.get_embedding(slide_info['content'])
            
            self.collection.add(
                documents=[slide_info['content']],
                embeddings=[embedding],
                ids=[f"{Path(pptx_path).stem}_slide_{slide_info['slide_number']}"],
                metadatas=[{
                    "source": pptx_path,
                    "slide_number": slide_info['slide_number'],
                    "filename": slide_info['filename'],
                    "type": "slide"
                }],
            )

    def extract_pdf_content(self, pdf_path: str):
        """Extract text content from PDF file."""
        content = ""
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages, start=1):
                content += f"\n--- Page {page_num} ---\n"
                content += page.extract_text()
        return content

    def analyze_practice_test(self, test_pdf_path: str, flagged_questions: list = None):
        """
        Analyze practice test and identify topics to review.
        
        Args:
            test_pdf_path: Path to the practice test PDF
            flagged_questions: List of question numbers that were wrong or need review
        
        Returns:
            Dictionary with analysis results
        """
        # Extract test content
        test_content = self.extract_pdf_content(test_pdf_path)
        
        # Use GPT to identify topics and questions
        analysis_prompt = f"""Analyze this practice test and extract:
1. All question topics/concepts being tested
2. If specific question numbers are flagged, focus on those questions

Practice Test Content:
{test_content}

{f"Focus on these question numbers: {', '.join(map(str, flagged_questions))}" if flagged_questions else "Analyze all questions"}

Return a structured list of:
- Question number
- Topic/concept being tested
- Key terms and keywords related to this question
"""
        
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are an expert at analyzing practice tests and identifying key topics and concepts."},
                {"role": "user", "content": analysis_prompt}
            ],
            temperature=0.3,
            max_tokens=2000
        )
        
        analysis = response.choices[0].message["content"]
        
        return {
            "test_analysis": analysis,
            "test_content": test_content
        }

    def find_relevant_slides(self, topics: str, n_results: int = 10):
        """Find slides relevant to specific topics."""
        # Get embedding for topics
        query_embedding = self.get_embedding(topics)
        
        # Search only slide content
        results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=n_results,
            where={"type": "slide"}
        )
        
        # Organize by file and slide number
        slides_by_file = {}
        for doc, metadata in zip(results["documents"][0], results["metadatas"][0]):
            filename = metadata.get("filename", "Unknown")
            slide_num = metadata.get("slide_number", "Unknown")
            
            if filename not in slides_by_file:
                slides_by_file[filename] = []
            
            slides_by_file[filename].append({
                "slide_number": slide_num,
                "content": doc,
                "source": metadata.get("source", "Unknown")
            })
        
        # Sort slides by number within each file
        for filename in slides_by_file:
            slides_by_file[filename].sort(key=lambda x: x["slide_number"])
        
        return slides_by_file

    def create_targeted_study_guide(self, test_pdf_path: str, flagged_questions: list = None):
        """
        Create a personalized study guide based on practice test performance.
        
        Args:
            test_pdf_path: Path to practice test PDF
            flagged_questions: List of question numbers to focus on (wrong/flagged)
        
        Returns:
            Dictionary with study guide and slide recommendations
        """
        print("📝 Analyzing practice test...")
        # Analyze the practice test
        test_analysis = self.analyze_practice_test(test_pdf_path, flagged_questions)
        
        print("🔍 Finding relevant slides...")
        # Find relevant slides based on the analysis
        relevant_slides = self.find_relevant_slides(test_analysis["test_analysis"], n_results=15)
        
        print("📚 Generating study guide...")
        # Create comprehensive study guide
        study_guide_prompt = f"""Based on this practice test analysis and the relevant course material, create a comprehensive study guide.

Practice Test Analysis:
{test_analysis["test_analysis"]}

Relevant Course Material:
{self._format_slides_for_prompt(relevant_slides)}

Create a study guide that:
1. Explains each concept that appeared in flagged questions
2. Provides clear definitions and examples
3. Highlights common mistakes or misconceptions
4. Includes memory aids or mnemonics where helpful
5. Suggests practice problems or ways to apply the concepts

Format the guide with clear sections and bullet points."""

        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are an expert tutor creating personalized study guides."},
                {"role": "user", "content": study_guide_prompt}
            ],
            temperature=0.5,
            max_tokens=3000
        )
        
        study_guide = response.choices[0].message["content"]
        
        return {
            "study_guide": study_guide,
            "slides_to_review": relevant_slides,
            "test_analysis": test_analysis["test_analysis"]
        }

    def _format_slides_for_prompt(self, slides_by_file: dict, max_slides: int = 10):
        """Format slide content for GPT prompt."""
        formatted = ""
        count = 0
        
        for filename, slides in slides_by_file.items():
            formatted += f"\n\nFrom {filename}:\n"
            for slide in slides[:max_slides]:
                if count >= max_slides:
                    break
                formatted += f"\nSlide {slide['slide_number']}:\n{slide['content'][:500]}...\n"
                count += 1
        
        return formatted

    def process_transcription(self, file_path):
        """Process a transcription file and add it to the vector database."""
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        # Split content into chunks
        chunks = self.text_splitter.split_text(content)

        # Generate embeddings and add to ChromaDB
        for i, chunk in enumerate(tqdm(chunks, desc="Processing chunks")):
            embedding = self.get_embedding(chunk)

            self.collection.add(
                documents=[chunk],
                embeddings=[embedding],
                ids=[f"{Path(file_path).stem}_{i}"],
                metadatas=[{"source": file_path, "chunk_id": i, "type": "text"}],
            )

    def query_knowledge_base(self, query, n_results=3):
        """Query the knowledge base with a natural language question."""
        # Generate query embedding
        query_embedding = self.get_embedding(query)

        # Search for relevant chunks
        results = self.collection.query(
            query_embeddings=[query_embedding], n_results=n_results
        )

        # Prepare context from results
        context = "\n\n".join(results["documents"][0])

        # Generate response using GPT
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {
                    "role": "system",
                    "content": """You are an AI study assistant for the AI Applications class. 
                Use the provided context to answer questions accurately and helpfully. 
                If the context doesn't contain enough information to answer fully, acknowledge this and suggest what additional information might be needed.""",
                },
                {"role": "user", "content": f"Context:\n{context}\n\nQuestion: {query}"},
            ],
            temperature=0.3,
            max_tokens=500,
        )

        return {
            "answer": response.choices[0].message["content"],
            "source_chunks": results["documents"][0],
            "metadata": results["metadatas"][0],
        }

    def generate_quiz(self, topic):
        """Generate a quiz based on the content in the knowledge base."""
        # First, retrieve relevant content about the topic
        query_embedding = self.get_embedding(topic)
        results = self.collection.query(
            query_embeddings=[query_embedding], n_results=5
        )

        context = "\n\n".join(results["documents"][0])

        # Generate quiz using GPT
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {
                    "role": "system",
                    "content": """Generate a quiz based on the provided content. 
                Create 3 questions of varying difficulty (easy, medium, hard).
                Each question should test understanding rather than mere recall.
                Include answers and explanations.""",
                },
                {"role": "user", "content": f"Content to base quiz on:\n{context}"},
            ],
            temperature=0.7,
            max_tokens=1000,
        )

        return response.choices[0].message["content"]

    def create_study_guide(self, topic):
        """Create a focused study guide on a specific topic."""
        query_embedding = self.get_embedding(topic)
        results = self.collection.query(
            query_embeddings=[query_embedding], n_results=5
        )

        context = "\n\n".join(results["documents"][0])

        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {
                    "role": "system",
                    "content": """Create a comprehensive study guide based on the provided content.
                Include:
                1. Key concepts and definitions
                2. Important relationships and connections
                3. Common misconceptions
                4. Practice problems or examples
                5. Summary of main points""",
                },
                {"role": "user", "content": f"Content to base study guide on:\n{context}"},
            ],
            temperature=0.5,
            max_tokens=1500,
        )

        return response.choices[0].message["content"]

    def concept_map(self, topic):
        """Generate a textual concept map showing relationships between ideas."""
        query_embedding = self.get_embedding(topic)
        results = self.collection.query(
            query_embeddings=[query_embedding], n_results=4
        )

        context = "\n\n".join(results["documents"][0])

        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {
                    "role": "system",
                    "content": """Create a textual concept map showing relationships between key ideas.
                Use simple ASCII/Unicode characters to show connections.
                Format should be clear and readable in a monospace font.""",
                },
                {"role": "user", "content": f"Content to map:\n{context}"},
            ],
            temperature=0.4,
            max_tokens=1000,
        )

        return response.choices[0].message["content"]
