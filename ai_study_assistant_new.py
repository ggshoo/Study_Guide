import os
from pathlib import Path
import chromadb
from chromadb.config import Settings
import openai
import json
from langchain.text_splitter import RecursiveCharacterTextSplitter
from tqdm import tqdm
from dotenv import load_dotenv
import PyPDF2
from pptx import Presentation
import re

# Load environment variables
load_dotenv()


class AIStudyAssistant:
    def __init__(self, persist_directory="db"):
        """Initialize the AI Study Assistant with a ChromaDB instance."""
        self.client = chromadb.PersistentClient(path=persist_directory)
        self.collection = self.client.get_or_create_collection(
            name="lecture_content",
            metadata={"hnsw:space": "cosine"}
        )

        # Configure API key if present; UI may set it later at runtime.
        env_key = os.getenv("OPENAI_API_KEY")
        if env_key:
            openai.api_key = env_key

        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        self.GEN_MODEL = os.getenv("GEN_MODEL", "gpt-4o-mini")
        # Embedding cache for performance (text hash -> embedding vector)
        self._embedding_cache = {}

    def get_embeddings_batch(self, texts, use_cache=True):
        """Batch embedding with optional caching to avoid redundant API calls."""
        import hashlib
        
        if not use_cache:
            resp = openai.Embedding.create(model="text-embedding-3-small", input=texts)
            return [d["embedding"] for d in resp["data"]]
        
        # Check cache and only embed uncached texts
        results = [None] * len(texts)
        texts_to_embed = []
        indices_to_embed = []
        
        for i, text in enumerate(texts):
            text_hash = hashlib.sha256(text.encode('utf-8')).hexdigest()
            if text_hash in self._embedding_cache:
                results[i] = self._embedding_cache[text_hash]
            else:
                texts_to_embed.append(text)
                indices_to_embed.append(i)
        
        # Embed uncached texts
        if texts_to_embed:
            resp = openai.Embedding.create(model="text-embedding-3-small", input=texts_to_embed)
            embeddings = [d["embedding"] for d in resp["data"]]
            
            # Cache and populate results
            for idx, text, emb in zip(indices_to_embed, texts_to_embed, embeddings):
                text_hash = hashlib.sha256(text.encode('utf-8')).hexdigest()
                self._embedding_cache[text_hash] = emb
                results[idx] = emb
        
        return results

    def get_embedding(self, text: str):
        """Call OpenAI Embeddings API directly and return the vector."""
        resp = openai.Embedding.create(model="text-embedding-3-small", input=text)
        return resp["data"][0]["embedding"]

    def extract_pptx_content(self, pptx_path: str, original_filename: str | None = None):
        """Extract text content from PowerPoint file with slide numbers.
        original_filename overrides the tmp path name for user display.
        """
        prs = Presentation(pptx_path)
        slides_content = []
        display_name = original_filename or Path(pptx_path).name
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = f"\n--- Slide {slide_num} ---\n"
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            
            # Extract notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    slide_text += f"\nNotes: {notes_text}\n"
            
            slides_content.append({
                'slide_number': slide_num,
                'content': slide_text,
                'filename': display_name
            })
        
        return slides_content

    def process_pptx(self, pptx_path: str, original_filename: str | None = None):
        """Process PowerPoint file and add to vector database with slide tracking."""
        slides = self.extract_pptx_content(pptx_path, original_filename=original_filename)
        
        # Batch slides to reduce API calls
        BATCH_SIZE = 24
        for start in tqdm(range(0, len(slides), BATCH_SIZE), desc=f"Processing {Path(pptx_path).name}"):
            batch = slides[start:start + BATCH_SIZE]
            batch_texts = [s['content'] for s in batch]
            embeddings = self.get_embeddings_batch(batch_texts)
            ids = [f"{Path(pptx_path).stem}_slide_{s['slide_number']}" for s in batch]
            metadatas = [{
                "source": pptx_path,
                "slide_number": s['slide_number'],
                "filename": s['filename'],  # use original display name if provided
                "type": "slide"
            } for s in batch]

            self.collection.add(
                documents=batch_texts,
                embeddings=embeddings,
                ids=ids,
                metadatas=metadatas,
            )

    def extract_pdf_content(self, pdf_path: str):
        """Extract text content from PDF."""
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
        return text
    
    def extract_questions_and_answers(self, test_path: str):
        """Extract questions, correct answers, and user's answers from a practice test.
        
        Returns:
            dict with 'questions' (dict of q_num -> question_text),
            'correct_answers' (dict of q_num -> answer),
            'user_answers' (dict of q_num -> answer if found)
        """
        ext = Path(test_path).suffix.lower()
        if ext == ".pdf":
            test_content = self.extract_pdf_content(test_path)
        elif ext == ".pptx":
            slides = self.extract_pptx_content(test_path)
            test_content = "\n".join([s["content"] for s in slides])
        else:
            raise ValueError(f"Unsupported file type: {ext}")
        
        extraction_prompt = f"""Extract all questions, correct answers, AND the user's answers from this practice test.

Practice Test Content:
{test_content}

Return the data in this exact JSON format:
{{
    "questions": {{
        "1": "question text here",
        "2": "question text here",
        ...
    }},
    "correct_answers": {{
        "1": "A",
        "2": "B",
        ...
    }},
    "user_answers": {{
        "1": "B",
        "2": "B",
        ...
    }}
}}

IMPORTANT Instructions:
- Question numbers should be integers as strings
- For multiple choice, return just the letter (A, B, C, D, etc.)
- Look for any written answers, circled answers, or marked selections on the test
- If you find the answer key, extract it into correct_answers
- If you find user's marked/written answers, extract them into user_answers
- If no answer key is found, return an empty correct_answers dict
- If no user answers are found, return an empty user_answers dict
- Include ALL questions you can identify
"""
        
        response = openai.ChatCompletion.create(
            model=self.GEN_MODEL,
            messages=[
                {"role": "system", "content": "You are an expert at extracting structured information from practice tests. Always return valid JSON with ALL questions."},
                {"role": "user", "content": extraction_prompt}
            ],
            temperature=0.1,
            max_tokens=3000  # Increased to handle more questions
        )
        
        import json
        import re
        
        content = response.choices[0].message["content"]
        # Extract JSON from markdown code blocks if present
        json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', content, re.DOTALL)
        if json_match:
            content = json_match.group(1)
        
        try:
            result = json.loads(content)
            # Ensure all expected keys exist
            if 'questions' not in result:
                result['questions'] = {}
            if 'correct_answers' not in result:
                result['correct_answers'] = {}
            if 'user_answers' not in result:
                result['user_answers'] = {}
            return result
        except:
            # Fallback if parsing fails
            return {"questions": {}, "correct_answers": {}, "user_answers": {}}

    def analyze_practice_test(self, test_path: str, flagged_questions: list = None):
        """Analyze practice test (PDF or PPTX) and identify topics to review.

        Args:
            test_path: Path to the practice test file (.pdf or .pptx)
            flagged_questions: Optional list of question numbers to emphasize

        Returns:
            dict with keys: test_analysis (str), test_content (str)
        """
        ext = Path(test_path).suffix.lower()
        if ext == ".pdf":
            test_content = self.extract_pdf_content(test_path)
        elif ext == ".pptx":
            # Treat each slide as potential question/context block
            slides = self.extract_pptx_content(test_path)
            test_content = "\n".join([s["content"] for s in slides])
        else:
            raise ValueError(f"Unsupported practice test file type: {ext}")

        flagged_clause = (
            f"Focus on these question numbers: {', '.join(map(str, flagged_questions))}" if flagged_questions else "Analyze all questions"
        )

        analysis_prompt = f"""Analyze this practice test and extract detailed information for each question.

Practice Test Content:
{test_content}

{flagged_clause}

For each question (especially flagged/incorrect ones), provide:
1. Question number (if identifiable)
2. Main topic/concept being tested
3. Specific sub-topics or skills required
4. Key terms, formulas, or concepts needed to answer correctly
5. Common misconceptions or mistakes for this type of question

Format as a structured list with clear question-by-question breakdown."""

        response = openai.ChatCompletion.create(
            model=self.GEN_MODEL,
            messages=[
                {"role": "system", "content": "You are an expert at analyzing practice tests and identifying key topics and concepts."},
                {"role": "user", "content": analysis_prompt}
            ],
            temperature=0.3,
            max_tokens=1200
        )

        analysis = response.choices[0].message["content"]
        return {"test_analysis": analysis, "test_content": test_content}

    def find_relevant_slides(self, topics: str, n_results: int = 10):
        """Find slides relevant to specific topics."""
        # Get embedding for topics
        query_embedding = self.get_embedding(topics)
        
        # Search only slide content
        results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=n_results,
            where={"type": "slide"}
        )
        
        # Organize by file and slide number
        slides_by_file = {}
        for doc, metadata in zip(results["documents"][0], results["metadatas"][0]):
            filename = metadata.get("filename", "Unknown")
            slide_num = metadata.get("slide_number", "Unknown")
            
            if filename not in slides_by_file:
                slides_by_file[filename] = []
            
            slides_by_file[filename].append({
                "slide_number": slide_num,
                "content": doc,
                "source": metadata.get("source", "Unknown")
            })
        
        # Sort slides by number within each file
        for filename in slides_by_file:
            slides_by_file[filename].sort(key=lambda x: x["slide_number"])
        
        return slides_by_file

    def create_targeted_study_guide(self, test_path: str, flagged_questions: list = None, fast_mode: bool = False):
        """
        Create a personalized study guide based on practice test performance.
        
        Args:
            test_path: Path to practice test file (.pdf or .pptx)
            flagged_questions: List of question numbers to focus on (wrong/flagged)
            fast_mode: If True, use lighter processing for faster results
        
        Returns:
            Dictionary with study guide and slide recommendations
        """
        print("📝 Analyzing practice test...")
        # Analyze the practice test
        test_analysis = self.analyze_practice_test(test_path, flagged_questions)

        # Extract questions for explicit mapping
        questions_data = self.extract_questions_and_answers(test_path)
        questions = questions_data.get('questions', {})
        
        # VALIDATION: Ensure all questions extracted
        print(f"✓ Extracted {len(questions)} questions from test")
        if len(questions) == 0:
            print("⚠️ WARNING: No questions extracted from test!")
            return {
                "study_guide": "ERROR: Could not extract questions from the test. Please ensure the PDF contains readable text.",
                "question_slides_map": {},
                "test_analysis": "No questions found"
            }

        print("🔍 Matching each question to slides...")
        # For each question, find top matching slides
        question_slides_map = {}
        n_per_question = 3 if fast_mode else 5
        
        # Process ALL questions, not just flagged ones
        questions_to_process = questions.items()
        if flagged_questions:
            # If specific questions flagged, prioritize those but include all
            print(f"  Prioritizing {len(flagged_questions)} flagged questions")
        
        for q_num, q_text in questions_to_process:
            # Use embeddings to find relevant slides for each question
            slides = self.find_relevant_slides(q_text, n_results=n_per_question)
            question_slides_map[q_num] = slides
            
        print(f"✓ Mapped all {len(question_slides_map)} questions to slides")

        print("📚 Generating study guide...")
        # Adjust generation parameters for fast mode
        max_tokens = 3500 if not fast_mode else 2000  # Increased token limit
        temperature = 0.4 if fast_mode else 0.5

        # If too many questions, batch them to avoid truncation
        question_items = list(questions.items())
        batch_size = 10 if not fast_mode else 15
        study_guide_parts = []
        
        print(f"  Processing {len(question_items)} questions in {(len(question_items) + batch_size - 1) // batch_size} batches")
        
        for batch_idx, i in enumerate(range(0, len(question_items), batch_size)):
            batch_questions = dict(question_items[i:i+batch_size])
            batch_slides_map = {q: question_slides_map[q] for q in batch_questions}
            q_nums = sorted([int(q) for q in batch_questions.keys()])
            print(f"  Batch {batch_idx + 1}: Questions {q_nums[0]}-{q_nums[-1]}")
            
            batch_prompt = f"""Based on this practice test analysis and the relevant course material, create a {"concise" if fast_mode else "comprehensive"} study guide for the following questions:

Practice Test Analysis:
{test_analysis["test_analysis"]}

Relevant Course Material (from PowerPoint slides):
{self._format_slides_for_prompt(batch_slides_map)}

Questions:
{json.dumps(batch_questions, indent=2)}

CRITICAL: You MUST cover ALL {len(batch_questions)} questions listed above. For EACH question:

1. **Slide Mapping Explanation**: Clearly explain WHY each recommended slide is relevant to this question
   - Which specific medical concept does the slide cover?
   - How does it relate to the question being asked?

2. **Multiple Choice Analysis** (if applicable):
   - For each answer option (A, B, C, D, etc.), explain:
     * Why it's correct/incorrect
     * Which slide(s) contain information about this option
     * Key concepts to understand for ruling in/out this option

3. **Medical Concept Review**:
   - Explain the medical concept in detail with examples from the slides
   - Highlight what was likely misunderstood
   - Common clinical pitfalls or misconceptions

4. **Study Approach**:
   - Step-by-step guidance on how to approach similar questions
   - Key decision points or diagnostic reasoning

Format for EACH question:
### Question [#]: [Brief topic]
**Recommended Slides:** [List with explanations]
**Answer Analysis:** [For MC: cover each option with slide references]
**Key Concepts:** [Medical content review]
**Study Tips:** [How to approach similar questions]

Be medically accurate and thorough."""

            response = openai.ChatCompletion.create(
                model=self.GEN_MODEL,
                messages=[
                    {"role": "system", "content": "You are a medical education expert creating accurate study guides for medical students. Always cover ALL questions provided."},
                    {"role": "user", "content": batch_prompt}
                ],
                temperature=temperature,
                max_tokens=max_tokens
            )
            study_guide_parts.append(response.choices[0].message["content"])

        study_guide = "\n\n---\n\n".join(study_guide_parts)
        
        print(f"✓ Generated study guide with {len(study_guide_parts)} sections")

        return {
            "study_guide": study_guide,
            "question_slides_map": question_slides_map,
            "total_questions": len(questions),
            "questions_mapped": len(question_slides_map),
            "test_analysis": test_analysis["test_analysis"]
        }

    def _format_slides_for_prompt(self, slides_by_file: dict, max_slides: int = 10):
        """Format slide content for GPT prompt.
        
        Args:
            slides_by_file: Can be either:
                - dict of filename -> list of slides (old format)
                - dict of question_num -> dict of filename -> list of slides (new format)
        """
        formatted = ""
        count = 0
        
        # Check if this is the new nested format (question -> filename -> slides)
        first_key = next(iter(slides_by_file.keys())) if slides_by_file else None
        if first_key and isinstance(slides_by_file[first_key], dict):
            # New format: flatten question-level nesting
            file_slides_map = {}
            for q_num, files in slides_by_file.items():
                for filename, slides in files.items():
                    if filename not in file_slides_map:
                        file_slides_map[filename] = []
                    file_slides_map[filename].extend(slides)
            slides_by_file = file_slides_map
        
        # Now format as filename -> slides
        for filename, slides in slides_by_file.items():
            if count >= max_slides:
                break
            formatted += f"\n\nFrom {filename}:\n"
            for slide in slides[:max_slides]:
                if count >= max_slides:
                    break
                formatted += f"\nSlide {slide['slide_number']}:\n{slide['content'][:500]}...\n"
                count += 1
        
        return formatted

    def process_transcription(self, file_path, original_filename: str | None = None):
        """Process a transcription file and add it to the vector database.
        original_filename is stored for display instead of tmp path.
        """
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        # Split content into chunks
        chunks = self.text_splitter.split_text(content)

        # Generate embeddings and add to ChromaDB
        for i, chunk in enumerate(tqdm(chunks, desc="Processing chunks")):
            embedding = self.get_embedding(chunk)

            self.collection.add(
                documents=[chunk],
                embeddings=[embedding],
                ids=[f"{Path(file_path).stem}_{i}"],
                metadatas=[{
                    "source": file_path,
                    "chunk_id": i,
                    "type": "text",
                    "filename": original_filename or Path(file_path).name
                }],
            )

    def query_knowledge_base(self, query, n_results=3):
        """Query the knowledge base with a natural language question."""
        # Generate query embedding
        query_embedding = self.get_embedding(query)

        # Search for relevant chunks
        results = self.collection.query(
            query_embeddings=[query_embedding], n_results=n_results
        )

        # Prepare context from results
        context = "\n\n".join(results["documents"][0])

        # Generate response using GPT
        response = openai.ChatCompletion.create(
            model=self.GEN_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": """You are an AI study assistant for the AI Applications class. 
                Use the provided context to answer questions accurately and helpfully. 
                If the context doesn't contain enough information to answer fully, acknowledge this and suggest what additional information might be needed.""",
                },
                {"role": "user", "content": f"Context:\n{context}\n\nQuestion: {query}"},
            ],
            temperature=0.3,
            max_tokens=500,
        )

        return {
            "answer": response.choices[0].message["content"],
            "source_chunks": results["documents"][0],
            "metadata": results["metadatas"][0],
        }

    def generate_quiz(self, topic):
        """Generate a quiz based on the content in the knowledge base."""
        # First, retrieve relevant content about the topic
        query_embedding = self.get_embedding(topic)
        results = self.collection.query(
            query_embeddings=[query_embedding], n_results=5
        )

        context = "\n\n".join(results["documents"][0])

        # Generate quiz using GPT
        response = openai.ChatCompletion.create(
            model=self.GEN_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": """Generate a quiz based on the provided content. 
                Create 3 questions of varying difficulty (easy, medium, hard).
                Each question should test understanding rather than mere recall.
                Include answers and explanations.""",
                },
                {"role": "user", "content": f"Content to base quiz on:\n{context}"},
            ],
            temperature=0.7,
            max_tokens=1000,
        )

        return response.choices[0].message["content"]

    def create_study_guide(self, topic):
        """Create a focused study guide on a specific topic."""
        query_embedding = self.get_embedding(topic)
        results = self.collection.query(
            query_embeddings=[query_embedding], n_results=5
        )

        context = "\n\n".join(results["documents"][0])

        response = openai.ChatCompletion.create(
            model=self.GEN_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": """Create a comprehensive study guide based on the provided content.
                Include:
                1. Key concepts and definitions
                2. Important relationships and connections
                3. Common misconceptions
                4. Practice problems or examples
                5. Summary of main points""",
                },
                {"role": "user", "content": f"Content to base study guide on:\n{context}"},
            ],
            temperature=0.5,
            max_tokens=1500,
        )

        return response.choices[0].message["content"]

    def concept_map(self, topic):
        """Generate a textual concept map showing relationships between ideas."""
        query_embedding = self.get_embedding(topic)
        results = self.collection.query(
            query_embeddings=[query_embedding], n_results=4
        )

        context = "\n\n".join(results["documents"][0])

        response = openai.ChatCompletion.create(
            model=self.GEN_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": """Create a textual concept map showing relationships between key ideas.
                Use simple ASCII/Unicode characters to show connections.
                Format should be clear and readable in a monospace font.""",
                },
                {"role": "user", "content": f"Content to map:\n{context}"},
            ],
            temperature=0.4,
            max_tokens=1000,
        )

        return response.choices[0].message["content"]
